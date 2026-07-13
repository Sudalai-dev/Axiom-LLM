"""
Ingestion Module — Axiom Dataset Generation Platform.

Provides a pipeline framework for ingesting company knowledge from various
document formats (Markdown, PDF, DOCX, PPTX, CSV, XLSX, code, images,
meeting notes, specifications, emails) into structured dataset records.

This module defines the ingestion interface and a text-based implementation.
Binary format extractors (PDF, DOCX, PPTX, XLSX) require optional dependencies
and will gracefully degrade to metadata-only records if unavailable.
"""

import hashlib
import os
import re
import uuid
from datetime import datetime
from typing import Any, Dict, List, Optional


# Supported extensions grouped by category
FORMAT_REGISTRY = {
    "markdown": [".md", ".markdown", ".mdx"],
    "code": [".py", ".js", ".ts", ".java", ".go", ".rs", ".cs", ".cpp", ".c", ".h",
             ".yaml", ".yml", ".json", ".toml", ".xml", ".html", ".css", ".sql",
             ".sh", ".bash", ".ps1", ".bat", ".dockerfile", ".tf", ".hcl"],
    "csv": [".csv", ".tsv"],
    "plain_text": [".txt", ".log", ".ini", ".cfg", ".conf", ".env"],
    "pdf": [".pdf"],
    "docx": [".docx"],
    "pptx": [".pptx"],
    "xlsx": [".xlsx", ".xls"],
    "image": [".png", ".jpg", ".jpeg", ".gif", ".svg", ".bmp", ".webp"],
    "email": [".eml", ".msg"],
}


def _detect_format(file_path: str) -> str:
    """Detect the document format category from the file extension."""
    ext = os.path.splitext(file_path)[1].lower()
    for category, extensions in FORMAT_REGISTRY.items():
        if ext in extensions:
            return category
    return "unknown"


def _compute_file_hash(file_path: str) -> str:
    """Compute SHA-256 hash of a file for deduplication."""
    h = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def _extract_text_content(file_path: str, format_category: str) -> str:
    """Extract text content from a file based on its format category.
    Returns raw text for text-based formats. Returns empty string for
    binary formats that require optional dependencies."""
    if format_category in ("markdown", "code", "csv", "plain_text"):
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        except Exception:
            return ""

    if format_category == "pdf":
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(file_path)
            text = "\n".join(page.get_text() for page in doc)
            doc.close()
            return text
        except ImportError:
            return ""
        except Exception:
            return ""

    if format_category == "docx":
        try:
            import docx
            doc = docx.Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs)
        except ImportError:
            return ""
        except Exception:
            return ""

    if format_category == "pptx":
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text)
            return "\n".join(texts)
        except ImportError:
            return ""
        except Exception:
            return ""

    if format_category == "xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            texts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    texts.append(" ".join(str(c) for c in row if c is not None))
            wb.close()
            return "\n".join(texts)
        except ImportError:
            return ""
        except Exception:
            return ""

    return ""


def _extract_entities(text: str) -> List[str]:
    """Simple keyword-based entity extraction from text content."""
    entities = set()

    # Technology keywords
    tech_keywords = {
        "MQTT", "OPC-UA", "Modbus", "BACnet", "Kafka", "RabbitMQ", "NATS",
        "PostgreSQL", "MySQL", "MongoDB", "Redis", "Elasticsearch",
        "Kubernetes", "Docker", "Terraform", "AWS", "Azure", "GCP",
        "React", "Angular", "Vue", "FastAPI", "Spring Boot", "Django",
        "HIPAA", "PCI-DSS", "GDPR", "SOC 2", "ISO 27001", "OWASP",
        "ISA-95", "IEC 62443", "NIST", "FHIR", "HL7",
        "PLC", "SCADA", "HMI", "DCS", "VFD", "RTU",
    }
    for keyword in tech_keywords:
        if keyword.lower() in text.lower():
            entities.add(keyword)

    # Standard references (e.g., ISO 9001, IEC 61850)
    std_matches = re.findall(r"\b(?:ISO|IEC|IEEE|NIST|API|ASME|ANSI)\s*\d+(?:\.\d+)?(?:-\d+)?\b", text)
    entities.update(std_matches)

    return sorted(entities)


def _classify_document(text: str, file_path: str) -> Dict[str, str]:
    """Heuristic classification of document domain, industry, and category."""
    lowered = text.lower()
    filename = os.path.basename(file_path).lower()

    # Domain detection
    domain = "software_engineering"
    domain_scores = {
        "industrial_iot": ["mqtt", "opc-ua", "modbus", "scada", "plc", "sensor", "telemetry"],
        "mechanical": ["pump", "compressor", "bearing", "motor", "hvac", "valve", "conveyor"],
        "electrical": ["vfd", "transformer", "switchgear", "relay", "power distribution"],
        "cloud": ["aws", "azure", "gcp", "kubernetes", "terraform", "serverless"],
        "database": ["postgresql", "mysql", "schema", "index", "query", "migration"],
        "security": ["vulnerability", "encryption", "firewall", "owasp", "penetration"],
        "backend": ["api", "endpoint", "middleware", "rest", "graphql"],
        "frontend": ["react", "angular", "vue", "css", "responsive", "component"],
    }
    best_score = 0
    for d, keywords in domain_scores.items():
        score = sum(1 for k in keywords if k in lowered)
        if score > best_score:
            best_score = score
            domain = d

    # Industry detection
    industry = "general"
    industry_scores = {
        "manufacturing": ["factory", "production", "assembly", "oee", "batch", "quality"],
        "healthcare": ["patient", "clinical", "hipaa", "ehr", "fhir", "medical"],
        "energy": ["grid", "solar", "wind", "turbine", "generation", "utility"],
        "finance": ["ledger", "transaction", "payment", "banking", "fraud", "compliance"],
        "oil_and_gas": ["pipeline", "wellhead", "upstream", "downstream", "refinery"],
    }
    best_score = 0
    for ind, keywords in industry_scores.items():
        score = sum(1 for k in keywords if k in lowered)
        if score > best_score:
            best_score = score
            industry = ind

    # Document category detection
    category = "specification"
    category_keywords = {
        "architecture_diagram": ["architecture", "diagram", "component"],
        "meeting_notes": ["meeting", "minutes", "action items", "attendees"],
        "lessons_learned": ["lesson", "retrospective", "what went well", "improvement"],
        "engineering_decision": ["decision", "adr", "alternative", "trade-off"],
        "incident_report": ["incident", "outage", "root cause", "postmortem"],
        "design_review": ["design review", "review comments", "approval"],
    }
    for cat, keywords in category_keywords.items():
        if any(k in lowered for k in keywords):
            category = cat
            break

    return {"domain": domain, "industry": industry, "category": category}


class IngestionPipeline:
    """Processes files from a directory into structured knowledge records."""

    def __init__(self, source_dir: str) -> None:
        self.source_dir = source_dir

    def scan(self) -> List[str]:
        """Scans the source directory and returns a list of supported file paths."""
        supported = set()
        for _, exts in FORMAT_REGISTRY.items():
            supported.update(exts)

        results = []
        for root, dirs, files in os.walk(self.source_dir):
            for f in files:
                ext = os.path.splitext(f)[1].lower()
                if ext in supported:
                    results.append(os.path.join(root, f))
        return sorted(results)

    def ingest(self, file_paths: Optional[List[str]] = None) -> List[Dict[str, Any]]:
        """Ingest files and produce structured knowledge records."""
        if file_paths is None:
            file_paths = self.scan()

        records = []
        for fp in file_paths:
            record = self.ingest_file(fp)
            if record:
                records.append(record)
        return records

    def ingest_file(self, file_path: str) -> Optional[Dict[str, Any]]:
        """Ingest a single file into a structured record."""
        if not os.path.exists(file_path):
            return None

        format_category = _detect_format(file_path)
        file_hash = _compute_file_hash(file_path)
        text_content = _extract_text_content(file_path, format_category)
        entities = _extract_entities(text_content) if text_content else []
        classification = _classify_document(text_content or "", file_path)

        now = datetime.utcnow().isoformat()
        stat = os.stat(file_path)

        return {
            "id": str(uuid.uuid4()),
            "version": "1.0",
            "source": "company_ingestion",
            "tags": [format_category, classification["category"]],
            "confidence": 0.80 if text_content else 0.50,
            "review_status": "pending_review",
            "created_date": datetime.fromtimestamp(stat.st_ctime).isoformat(),
            "updated_date": now,
            "approval": "pending",
            "domain": classification["domain"],
            "industry": classification["industry"],
            "intent": "company_knowledge",
            "complexity": "medium",
            "file_path": file_path,
            "file_name": os.path.basename(file_path),
            "file_hash": file_hash,
            "file_size_bytes": stat.st_size,
            "document_format": format_category,
            "document_category": classification["category"],
            "text_content_length": len(text_content),
            "extracted_entities": entities,
            "ingestion_status": "processed" if text_content else "metadata_only",
            "summary": text_content[:500] if text_content else f"Binary file: {os.path.basename(file_path)}",
        }
